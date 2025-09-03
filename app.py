import os
import streamlit as st
import openai
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import fitz  # PyMuPDF for PDFs
import docx
from pptx import Presentation

# --- Helper Functions ---
def load_pdf(file):
    text = ""
    with fitz.open(stream=file.read(), filetype="pdf") as doc:
        for page in doc:
            text += page.get_text("text") + "\n"
    return text

def load_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def load_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def load_txt(file):
    return file.read().decode("utf-8")

def chunk_text(text, chunk_size=500):
    words = text.split()
    return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def build_index(docs):
    chunks = []
    for doc_text in docs:
        chunks.extend(chunk_text(doc_text))
    vectorizer = TfidfVectorizer().fit(chunks)
    vectors = vectorizer.transform(chunks)
    return chunks, vectorizer, vectors

def retrieve(query, chunks, vectorizer, vectors, top_k=3):
    query_vec = vectorizer.transform([query])
    sims = cosine_similarity(query_vec, vectors).flatten()
    top_indices = sims.argsort()[-top_k:][::-1]
    return [chunks[i] for i in top_indices]

def ask_gpt(query, context):
    openai.api_key = os.getenv("OPENAI_API_KEY")
    prompt = f"Answer the following question using ONLY the provided context:\n\nContext:\n{context}\n\nQuestion: {query}\nAnswer:"
    
    response = openai.ChatCompletion.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=300,
        temperature=0.3,
    )
    return response["choices"][0]["message"]["content"]

# --- Streamlit App ---
st.set_page_config(page_title="Study Buddy Chatbot", layout="wide")
st.title("📚 Study Buddy Chatbot")
st.write("Upload your **textbooks, slides, and notes**, then ask me questions!")

uploaded_files = st.file_uploader(
    "Upload files (PDF, DOCX, PPTX, TXT)", type=["pdf", "docx", "pptx", "txt"], accept_multiple_files=True
)

if uploaded_files:
    st.success(f"Uploaded {len(uploaded_files)} files successfully!")

    docs = []
    for file in uploaded_files:
        if file.type == "application/pdf":
            docs.append(load_pdf(file))
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            docs.append(load_docx(file))
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            docs.append(load_pptx(file))
        elif file.type == "text/plain":
            docs.append(load_txt(file))

    # Build TF-IDF index
    st.info("Building knowledge index...")
    chunks, vectorizer, vectors = build_index(docs)
    st.success("Index built! You can now ask questions.")

    # Chat Interface
    query = st.text_input("Ask a question:")
    if query:
        with st.spinner("Thinking..."):
            retrieved = retrieve(query, chunks, vectorizer, vectors)
            context = "\n\n".join(retrieved)
            answer = ask_gpt(query, context)
        st.subheader("💡 Answer")
        st.write(answer)

        with st.expander("🔍 Retrieved Context"):
            st.write(context)
